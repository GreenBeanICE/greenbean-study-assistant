"""
PPT 解析器，用于解析 .pptx 格式文件，提取文本、标题和备注。
基于 python-pptx 库实现。

输出统一的 PageIndex 结构，与 PDFParser 保持一致。
"""
import io
from typing import List, Dict, Any
from pptx import Presentation
from app.utils.text_utils import clean_text

# ---- 常量和辅助函数 ----

def _is_title_placeholder(shape) -> bool:
    """检测 shape 是否为标题占位符（idx == 0）。"""
    if not shape.is_placeholder:
        return False
    return shape.placeholder_format.idx == 0


def _extract_slide_texts_and_headings(slide) -> tuple[List[str], List[Dict[str, Any]]]:
    """
    从单张 slide 中提取全部文本和标题列表。
    
    :return: (slide_texts, headings)
    """
    slide_texts: List[str] = []
    headings: List[Dict[str, Any]] = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            slide_texts.append(text)

            if _is_title_placeholder(shape):
                headings.append({"level": 1, "text": text})

    return slide_texts, headings


def _build_page_node(
    slide_num: int,
    slide_texts: List[str],
    headings: List[Dict[str, Any]],
    shapes_count: int,
) -> Dict[str, Any]:
    """构建单张 slide 对应的 PageIndex 节点。"""
    full_text = clean_text("\n".join(slide_texts))

    return {
        "page_number": slide_num,
        "content": full_text,
        "char_count": len(full_text),
        "parser_name": "PptParser",
        "parser_version": "1.0.0",
        "metadata": {
            "source_type": "ppt",
            "headings": headings,
            "shapes_count": shapes_count,
            "paragraphs_count": len(slide_texts),
        },
    }


class PptParser:
    """解析 .pptx 格式的 PowerPoint 文件，提取结构化文本内容。"""

    def parse(self, file_content: bytes) -> List[Dict[str, Any]]:
        """
        解析 PPT 字节流，逐页提取文本并构建 PageIndex 结构。
        
        每张幻灯片视为 1 个 PageIndex 节点。
        
        :param file_content: 前端上传文件的二进制数据
        :return: 包含每页文本和元数据的列表
        """
        prs = Presentation(io.BytesIO(file_content))
        parsed_pages: List[Dict[str, Any]] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts, headings = _extract_slide_texts_and_headings(slide)
            page_node = _build_page_node(
                slide_num,
                slide_texts,
                headings,
                shapes_count=len(slide.shapes),
            )
            parsed_pages.append(page_node)

        return parsed_pages
