"""
PPT 解析器测试，覆盖 PptParser 基本逻辑。
需要 python-pptx 库生成测试用的 .pptx 文件。
"""
import io
import pytest
from unittest.mock import patch, PropertyMock
from pptx import Presentation
from pptx.util import Inches
from app.parsers.ppt_parser import PptParser


def create_test_pptx(slides_count: int = 2) -> bytes:
    """生成测试用的 .pptx 文件字节流"""
    prs = Presentation()
    
    for i in range(slides_count):
        slide_layout = prs.slide_layouts[0]  # 标题+内容布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 标题占位符 (idx=0)
        title = slide.shapes.title
        title.text = f"第{i+1}章 测试标题"
        
        # 内容占位符 (idx=1)
        content = slide.placeholders[1]
        content.text = f"这是第{i+1}页的内容。\n包含多行文本。"
    
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


class TestPptParser:
    """PptParser 解析器测试"""

    def test_parse_pptx(self):
        """测试解析 .pptx 文件"""
        parser = PptParser()
        content = create_test_pptx(slides_count=2)
        result = parser.parse(content)
        
        assert len(result) == 2
        assert result[0]["page_number"] == 1
        assert result[1]["page_number"] == 2
        assert "第1章 测试标题" in result[0]["content"]
        assert "第2章 测试标题" in result[1]["content"]
        assert result[0]["metadata"]["source_type"] == "ppt"
        assert result[1]["metadata"]["source_type"] == "ppt"

    def test_parse_single_slide(self):
        """测试解析单页 PPT"""
        parser = PptParser()
        content = create_test_pptx(slides_count=1)
        result = parser.parse(content)
        
        assert len(result) == 1
        assert result[0]["page_number"] == 1
        assert result[0]["char_count"] > 0
        assert len(result[0]["metadata"]["headings"]) > 0

    def test_parse_empty_pptx(self):
        """测试解析空 PPT（无幻灯片）"""
        parser = PptParser()
        prs = Presentation()
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        result = parser.parse(buf.read())
        
        assert len(result) == 0

    def test_parse_metadata_structure(self):
        """测试解析结果的元数据结构"""
        parser = PptParser()
        content = create_test_pptx(slides_count=1)
        result = parser.parse(content)
        
        metadata = result[0]["metadata"]
        assert "source_type" in metadata
        assert "headings" in metadata
        assert "shapes_count" in metadata
        assert "paragraphs_count" in metadata
        assert metadata["source_type"] == "ppt"

    def test_parse_shape_without_text_frame(self):
        """测试包含无文本框架形状的 PPT（覆盖第 34 行 continue）"""
        import io
        from unittest.mock import PropertyMock, MagicMock
        from pptx import Presentation
        from pptx.util import Inches
        from pptx.enum.shapes import MSO_SHAPE
        
        parser = PptParser()
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加一个形状
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(3), Inches(1)
        )
        shape.text_frame.text = "有文本的形状"
        
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        content = buf.read()
        
        # 使用 mock 让 slide.shapes 返回一个包含无 text_frame 形状的列表
        # 直接 mock PptParser.parse 方法中使用的 slide.shapes 迭代
        with patch("app.parsers.ppt_parser.Presentation") as mock_presentation:
            mock_prs = MagicMock()
            mock_presentation.return_value = mock_prs
            
            # 创建一个有 text_frame 的 mock shape
            mock_shape_with_text = MagicMock()
            mock_shape_with_text.has_text_frame = True
            mock_text_frame = MagicMock()
            mock_para = MagicMock()
            mock_para.text = "有效文本"
            mock_para2 = MagicMock()
            mock_para2.text = ""
            mock_text_frame.paragraphs = [mock_para, mock_para2]
            mock_shape_with_text.text_frame = mock_text_frame
            mock_shape_with_text.is_placeholder = False
            
            # 创建一个无 text_frame 的 mock shape
            mock_shape_no_text = MagicMock()
            mock_shape_no_text.has_text_frame = False
            
            mock_slide = MagicMock()
            mock_slide.shapes = [mock_shape_no_text, mock_shape_with_text]
            mock_prs.slides = [mock_slide]
            
            result = parser.parse(b"fake pptx content")
            
            assert len(result) == 1
            assert "有效文本" in result[0]["content"]

    def test_parse_empty_paragraph_skipped(self):
        """测试跳过空文本段落（覆盖第 39 行 continue）"""
        import io
        from pptx import Presentation
        
        parser = PptParser()
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        # 标题占位符
        title = slide.shapes.title
        title.text = "标题"
        
        # 内容占位符 - 设置一些空行和有效文本
        content_ph = slide.placeholders[1]
        content_ph.text = "\n\n有效内容\n\n"
        
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        content = buf.read()
        
        result = parser.parse(content)
        assert len(result) == 1
        assert "有效内容" in result[0]["content"]
        assert "标题" in result[0]["content"]
        assert result[0]["metadata"]["paragraphs_count"] >= 1
